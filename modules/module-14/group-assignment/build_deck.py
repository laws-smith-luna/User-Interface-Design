"""
Build a visually polished PowerPoint deck for the SWE 632 final group presentation.

Run:  python build_deck.py
Output: final-presentation.pptx in the same folder

Design language:
- Dark hero title slide
- Three colored section dividers (one per speaker)
- Clean content slides with a colored title bar and consistent footer
- Brand palette tied to Budget Buddy (greens/teals, slate neutrals)

Recording workflow: open in PowerPoint, use Slide Show > Record per slide.
Each speaker records their assigned slides. Demo slide gets an embedded MP4.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).parent / "final-presentation.pptx"

# -------- Brand palette --------
DARK     = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
DARK_2   = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
GREEN    = RGBColor(0x16, 0xA3, 0x4A)   # savings green
TEAL     = RGBColor(0x08, 0x91, 0xB2)   # accent teal
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # warning
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_300= RGBColor(0xCB, 0xD5, 0xE1)
SLATE_500= RGBColor(0x64, 0x74, 0x8B)
SLATE_700= RGBColor(0x33, 0x41, 0x55)

FONT = "Calibri"

# -------- Slide setup (16:9) --------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# -------- Helpers --------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=18, color=SLATE_700, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        # Bullet char + text
        bullet_run = p.add_run()
        bullet_run.text = "•  "
        bullet_run.font.name = font
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = GREEN
        bullet_run.font.bold = True

        run = p.add_run()
        run.text = b
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# -------- Slide builders --------
def title_slide():
    s = prs.slides.add_slide(BLANK)
    # Dark full-bleed background
    add_rect(s, 0, 0, SW, SH, DARK)
    # Accent vertical rule
    add_rect(s, Inches(0.9), Inches(2.4), Inches(0.18), Inches(2.7), GREEN)
    # Wordmark
    add_text(s, Inches(1.3), Inches(2.3), Inches(11), Inches(1.5),
             "Budget Buddy", size=80, bold=True, color=WHITE)
    # Tagline
    add_text(s, Inches(1.3), Inches(3.7), Inches(11), Inches(0.7),
             "Stay on pace, not just on budget.", size=28, color=SLATE_300)
    # Team line
    add_text(s, Inches(1.3), Inches(4.8), Inches(11), Inches(0.4),
             "Laws Smith   ·   Kevin Le   ·   Samana Hussain",
             size=18, color=SLATE_300)
    # Course
    add_text(s, Inches(1.3), Inches(5.3), Inches(11), Inches(0.4),
             "SWE 632 — User Interface Design and Development   ·   Spring 2026",
             size=14, color=SLATE_500)
    # Slide number badge bottom right
    add_text(s, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
             "01", size=12, color=SLATE_500, align=PP_ALIGN.RIGHT)
    set_notes(s,
        "[LAWS — 0:15]\n"
        "Hi, we're team Budget Buddy. I'm Laws, and with me are Kevin and Samana. "
        "Over the last semester we built a personal budgeting app focused on one question: "
        "am I on pace, or am I overspending? Here's what we built and how we got there.")
    return s


def section_divider(num, title, speaker, color, slide_num):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, color)
    # Big number
    add_text(s, Inches(0.9), Inches(1.3), Inches(6), Inches(2),
             num, size=180, bold=True, color=WHITE)
    # Title
    add_text(s, Inches(0.95), Inches(4.3), Inches(11), Inches(1),
             title, size=44, bold=True, color=WHITE)
    # Speaker line
    add_text(s, Inches(0.95), Inches(5.3), Inches(11), Inches(0.5),
             f"Presented by {speaker}", size=20, color=WHITE)
    # Footer slide number
    add_text(s, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
             f"{slide_num:02d}", size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    return s


def content_slide(title, bullets, notes, owner, slide_num,
                  bar_color=GREEN, hero_text=None):
    s = prs.slides.add_slide(BLANK)
    # Top accent bar
    add_rect(s, 0, 0, SW, Inches(0.9), bar_color)
    # Title in bar
    add_text(s, Inches(0.6), Inches(0.18), Inches(10), Inches(0.6),
             title, size=28, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # Owner badge top right
    add_text(s, Inches(11.3), Inches(0.18), Inches(1.9), Inches(0.6),
             owner, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Side accent rule (full height under bar)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(0.06), Inches(5.0), bar_color)

    # Optional hero callout above bullets
    body_y = Inches(1.2)
    if hero_text:
        add_text(s, Inches(0.85), Inches(1.2), Inches(11.8), Inches(0.7),
                 hero_text, size=20, bold=True, color=DARK)
        body_y = Inches(2.0)

    # Bullets
    add_bullets(s, Inches(0.85), body_y, Inches(11.8), Inches(5.2),
                bullets, size=18, color=SLATE_700)

    # Footer
    add_rect(s, 0, Inches(7.05), SW, Inches(0.45), SLATE_50)
    add_text(s, Inches(0.6), Inches(7.13), Inches(8), Inches(0.3),
             "Budget Buddy  ·  SWE 632 Final Presentation",
             size=10, color=SLATE_500)
    add_text(s, Inches(11.0), Inches(7.13), Inches(2.0), Inches(0.3),
             f"{slide_num:02d} / 14",
             size=10, color=SLATE_500, align=PP_ALIGN.RIGHT)

    set_notes(s, notes)
    return s


def thanks_slide(slide_num):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DARK)
    add_rect(s, Inches(0.9), Inches(2.4), Inches(0.18), Inches(2.0), GREEN)
    add_text(s, Inches(1.3), Inches(2.3), Inches(11), Inches(1.2),
             "Thank you.", size=72, bold=True, color=WHITE)
    add_text(s, Inches(1.3), Inches(3.8), Inches(11), Inches(0.6),
             "Questions?", size=28, color=SLATE_300)
    add_text(s, Inches(1.3), Inches(5.0), Inches(11), Inches(0.4),
             "Laws Smith   ·   Kevin Le   ·   Samana Hussain",
             size=18, color=SLATE_300)
    add_text(s, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
             f"{slide_num:02d}", size=12, color=SLATE_500, align=PP_ALIGN.RIGHT)
    set_notes(s, "[ALL — 0:15]\nThanks for watching. Happy to answer any questions.")
    return s


# ============================================================================
# DECK ASSEMBLY
# ============================================================================

# Slide 1 — Title
title_slide()

# Slide 2 — Section: The Project (Laws)
section_divider("01", "The Project", "Laws Smith", DARK_2, 2)

# Slide 3 — The Problem
content_slide(
    title="The Problem",
    hero_text="Most budgeting apps fail at the moment that matters.",
    bullets=[
        "Existing tools are scattered across disconnected places",
        "They demand too much manual categorization to stay accurate",
        "And they only update users after the money has already been spent",
        "The save-able window is the moment of decision, not after",
    ],
    notes=("[LAWS — 0:45]\n"
           "Most budgeting tools fall short the same way. They're scattered across disconnected "
           "places, they demand a lot of manual categorization, and they really only update you "
           "after the money has already been spent. The window where a budget can actually shape "
           "behavior is the moment of decision, not after. That's the window we set out to address."),
    owner="LAWS",
    slide_num=3,
    bar_color=DARK_2,
)

# Slide 4 — Our Solution
content_slide(
    title="Our Solution",
    hero_text="Per-category budgets translated into daily spending pace.",
    bullets=[
        "Monthly budgets for Food, Clothing, and other categories",
        "'$47 left, 2 days remaining, ~$23/day' — pace, not just balance",
        "Alerts surface early, before a category is already overspent",
        "Five sections: Dashboard · Budgets · Transactions · Alerts · Settings",
        "[INSERT hero screenshot of final Dashboard]",
    ],
    notes=("[LAWS — 1:00]\n"
           "Budget Buddy tracks per-category monthly budgets and translates remaining dollars into "
           "a daily spending allowance. Instead of '$47 left' with no context, you see '$47 left, "
           "2 days remaining, about $23 a day.' Alerts surface when a category is trending off-pace "
           "early enough to do something about it. The app is organized around five sections: "
           "Dashboard, Budgets, Transactions, Alerts, and Settings."),
    owner="LAWS",
    slide_num=4,
    bar_color=GREEN,
)

# Slide 5 — Project Origin & Initial Design
content_slide(
    title="Project Origin & Initial Design",
    hero_text="From a personal need to a testable scenario.",
    bullets=[
        "Started as Laws's M3 individual: 'Design a New Consumer Software Product'",
        "Years of using Quicken, Mint, YNAB, Rocket Money — none real-time AND all-inclusive",
        "Kevin and Samana joined the concept as a group project",
        "Formalized into the Sarah scenario via M5 hand-drawn wireframes",
        "[INSERT M5 hand-drawn wireframes — Dashboard, Clothing, Food]",
    ],
    notes=("[LAWS — 1:00]\n"
           "Budget Buddy started not from a formal needfinding study but from a personal need "
           "I articulated in my M3 individual assignment. I'd tried Quicken, Mint, YNAB, Rocket "
           "Money, and a rotation of bank apps over the years and none of them were both "
           "real-time and all-inclusive in the way I wanted. Kevin and Samana joined the concept "
           "as a group project. In M5 I formalized that into the Sarah scenario — a 28-year-old "
           "saving for a vacation who needs to make in-the-moment spending decisions. The "
           "hand-drawn wireframes covered the Dashboard, the Clothing and Food budget screens "
           "with donut charts, and the post-purchase push notification. That's the foundation "
           "every later iteration built on. Over to Kevin."),
    owner="LAWS",
    slide_num=5,
    bar_color=GREEN,
)

# Slide 6 — Section: First Budget Buddy Iterations (Kevin)
section_divider("02", "First Budget Buddy Iterations", "Kevin Le", TEAL, 6)

# Slide 7 — M6 Usability Study
content_slide(
    title="Iteration 2: M6 Usability Study",
    hero_text="Watching real users surfaced what reviews missed.",
    bullets=[
        "Converted Laws's M5 wireframes into interactive HTML (Wizard of Oz)",
        "Think-aloud study, 4 participants (A, B, C, D)",
        "5 usability issues identified — top 3 drove the M8 iteration",
        "No purchase preview · vague transaction labels · confusing budget creation",
        "[INSERT screenshot of M6 wireframe — Food Budget screen]",
    ],
    notes=("[KEVIN — 1:15]\n"
           "Our first group iteration on Budget Buddy was the M6 usability study. We converted "
           "Laws's hand-drawn wireframes into interactive HTML and ran a think-aloud study with "
           "four participants. Each of us ran at least one session. We identified five usability "
           "issues. The three most consequential were that users had no way to preview how a "
           "purchase would affect a budget — they could only see it after the fact — that "
           "transaction labels were too vague to verify against real spending, and that the "
           "budget-creation flow had several confusing moments. Those three findings drove the "
           "M8 iteration."),
    owner="KEVIN",
    slide_num=7,
    bar_color=TEAL,
)

# Slide 8 — M8 Interaction Iteration
content_slide(
    title="Iteration 3: M8 Interaction Iteration",
    hero_text="Three usability fixes from the M6 study.",
    bullets=[
        "Plan-a-Purchase preview on each budget detail screen",
        "Merchant-specific transaction labels (Chipotle, Starbucks, Trader Joe's)",
        "Clearer budget-creation flow: dollar amounts under percentages, full-width CTA",
        "[INSERT M8 before/after — Plan-a-Purchase preview]",
    ],
    notes=("[KEVIN — 1:00]\n"
           "M8 was a focused re-design pass that addressed three of the M6 findings. We added a "
           "Plan-a-Purchase card to each budget detail screen, so users can simulate a purchase "
           "and immediately see the projected balance. We replaced vague labels like 'Lunch' with "
           "merchant names. And we cleaned up the budget-creation flow — alert thresholds now "
           "show dollar amounts alongside percentages, and the New Budget button is now a clear "
           "primary CTA."),
    owner="KEVIN",
    slide_num=8,
    bar_color=TEAL,
)

# Slide 9 — M9 Interaction Critique
content_slide(
    title="Iteration 4: M9 Interaction Critique",
    hero_text="Self-critique against Site Design + Interaction + Preventing Error.",
    bullets=[
        "7 weaknesses identified, each tied to a cited principle",
        "Active bottom-nav highlight (Krug, 'You-are-here' marker)",
        "Confirmation step before committing a new budget (Norman, Ch. 5)",
        "Daily-pace context on budget cards (visibility of system status)",
        "[INSERT M9 before/after — confirmation flow]",
    ],
    notes=("[KEVIN — 1:00]\n"
           "M9 was a structured self-critique against three principle sets: Site Design, "
           "Interaction Techniques, and Preventing Error. We identified seven weaknesses and "
           "tied each one to a specific principle. The most consequential fixes: the bottom "
           "navigation now highlights the active section, the budget-creation flow has a "
           "confirmation step before committing, and the budget cards now show daily allowance "
           "and days remaining alongside the raw remaining-dollar number. Over to Samana."),
    owner="KEVIN",
    slide_num=9,
    bar_color=TEAL,
)

# Slide 10 — Section: Refinement & Demo (Samana)
section_divider("03", "Refinement & Demo", "Samana Hussain", GREEN, 10)

# Slide 11 — M10 + M13 Late Iterations
content_slide(
    title="Iterations 5 & 6: M10 + M13",
    hero_text="Sharpening interaction consistency, then visual hierarchy.",
    bullets=[
        "M10: dead-end Alerts/Settings tabs, contextual Quick Actions, back-button consistency",
        "M10: warning when creating a duplicate-category budget",
        "M13: critical alerts now visually dominant (hierarchy)",
        "M13: budget aggregate at top, threshold indicators on donuts, severity tags on alerts",
        "[INSERT M13 before/after — alert visual dominance]",
    ],
    notes=("[SAMANA — 1:15]\n"
           "M10 was a second round of interaction critique focused on consistency and error "
           "prevention. We fixed dead-end navigation tabs, replaced duplicate Quick Actions with "
           "contextual shortcuts, made back-button labels match the navigation hierarchy, and "
           "added a warning when creating a budget for a category that already had one. M13 was "
           "the visual-design pass — critiqued against Mullet & Sano and Tufte. Critical alerts "
           "are now visually dominant rather than blending in with status updates. We added a "
           "budget aggregate at the top of the Dashboard, threshold indicators on each donut, "
           "and explicit severity tags on alert cards so meaning survives without color."),
    owner="SAMANA",
    slide_num=11,
    bar_color=GREEN,
)

# Slide 12 — Demo
content_slide(
    title="Live Demo",
    hero_text="See it in action.",
    bullets=[
        "[Embed MP4 — recorded screen walkthrough]",
        "Task 1: Check your budget before a purchase (clothing + food)",
        "Task 2: Create a new vacation savings budget with an alert",
        "Task 3: Re-categorize a miscategorized Amazon transaction",
    ],
    notes=("[SAMANA — 1:15]\n"
           "TODO: Samana to record demo MP4 separately and embed on this slide. Narrate what "
           "the user sees and why it works. Three tasks: check budget before a purchase, create "
           "a new budget with the confirmation flow, and re-categorize a transaction."),
    owner="SAMANA",
    slide_num=12,
    bar_color=AMBER,
)

# Slide 13 — Reflection
content_slide(
    title="Reflection & Next Steps",
    hero_text="What we learned · what we'd build next.",
    bullets=[
        "Outside frameworks beat team intuition — every cycle",
        "Visual design is where perceived quality lives",
        "Next: real data via Plaid · predictive alerts · accessibility audit",
        "Bigger gap: dedicated needfinding before scaling the design",
    ],
    notes=("[SAMANA — 0:30]\n"
           "Two main learnings. Every iteration that started with 'what principle is this "
           "violating?' produced clearer fixes than ones that started with 'what should we "
           "change?' And visual design was where the biggest perceived-quality gains "
           "happened — M9 and M10 fixed real bugs, but M13 is what made the app feel like "
           "something we'd want to use. With more time: real data via Plaid, predictive alerts, "
           "an accessibility audit, and the dedicated needfinding study we never ran."),
    owner="SAMANA",
    slide_num=13,
    bar_color=GREEN,
)

# Slide 14 — Thanks
thanks_slide(14)


prs.save(OUT)
print(f"Wrote {OUT}")
print(f"Slides: {len(prs.slides)}")
