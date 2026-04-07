"""Build Angular Tech Talk PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1E, 0x1E, 0x2E)
BG_SLIDE = RGBColor(0x24, 0x24, 0x38)
ACCENT = RGBColor(0xDD, 0x33, 0x44)
ANGULAR_RED = RGBColor(0xDD, 0x00, 0x31)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
MID_GRAY = RGBColor(0x99, 0x99, 0xAA)
CODE_BG = RGBColor(0x18, 0x18, 0x28)
GREEN = RGBColor(0x66, 0xBB, 0x6A)
BLUE = RGBColor(0x64, 0xB5, 0xF6)
YELLOW = RGBColor(0xFF, 0xD5, 0x4F)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_code_box(slide, left, top, width, height, code, title=None):
    # Background shape for code
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Title above code if provided
    y_offset = top + 0.15
    if title:
        add_text_box(slide, left + 0.2, top + 0.1, width - 0.4, 0.4,
                     title, font_size=12, color=MID_GRAY, bold=True)
        y_offset = top + 0.45

    # Code text
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.3), Inches(y_offset),
        Inches(width - 0.6), Inches(height - (y_offset - top) - 0.1)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Cascadia Code"
        p.space_after = Pt(2)
    return tf


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_section_header(slide, number, title):
    set_slide_bg(slide, BG_DARK)
    # Section number
    add_text_box(slide, 0.8, 2.0, 11, 1.0,
                 f"0{number}" if number < 10 else str(number),
                 font_size=72, color=ANGULAR_RED, bold=True)
    # Section title
    add_text_box(slide, 0.8, 3.2, 11, 1.5, title,
                 font_size=44, color=WHITE, bold=True)


# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 1.5, 11, 1.5, "Angular", font_size=72, color=WHITE, bold=True)
add_text_box(slide, 0.8, 3.2, 11, 1.0,
             "A complete front-end framework for building web applications",
             font_size=28, color=LIGHT_GRAY)
# Red accent line
shape = slide.shapes.add_shape(1, Inches(0.8), Inches(3.0), Inches(3), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ANGULAR_RED
shape.line.fill.background()

add_text_box(slide, 0.8, 5.5, 11, 0.5, "SWE 632 Tech Talk",
             font_size=18, color=MID_GRAY)
add_text_box(slide, 0.8, 6.0, 11, 0.5, "Laws Smith",
             font_size=16, color=MID_GRAY)
add_notes(slide, "Hey everyone, my name's Laws. So today I want to talk about Angular. If you've ever used a web app that feels really snappy, like Gmail or Google Cloud, there's a good chance it was built with something like Angular. I'm going to walk you through what it is, show you what it looks like to actually build something with it, and then we'll talk about how it stacks up against other tools you might have heard of, like React or Vue.")


# ─── SLIDE 2: Section Header - Intro ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "Introduction & What is Angular")
add_notes(slide, "Alright, let's get into it.")


# ─── SLIDE 3: What is Angular ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "What is Angular?",
             font_size=36, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 1.3, 5.5, 5.0, [
    "Full front-end framework built and maintained by Google",
    "Written in TypeScript (not optional, it's baked in)",
    "Provides everything out of the box: routing, forms, HTTP client, testing",
    "Used for building single-page applications (SPAs)",
    "The page updates dynamically without full reloads",
], font_size=20)
# Timeline box
add_code_box(slide, 7.0, 1.3, 5.5, 4.5, """Timeline:

2010  AngularJS (1.x) launches
      One of the first major JS frameworks

2016  Angular 2 releases
      Complete rewrite in TypeScript
      Basically a brand new framework

2021  AngularJS reaches end-of-life
      No longer maintained

2024  Angular 19 (current)
      Major release every 6 months""", title="HISTORY")
add_notes(slide, """So Angular is a front-end framework that comes from Google. And I want to be clear about what "framework" means here, because it matters. Unlike something like React, which is really just a library for building UI pieces, Angular ships with everything. Routing, forms, an HTTP client, testing tools. You don't have to go shopping for plugins. It's all there.

Now one thing that trips people up is the name. There was actually a thing called AngularJS back in 2010, and it was huge. But in 2016, Google basically said "we're starting over" and rewrote the whole thing from scratch in TypeScript. Different language, different architecture, different everything. So when I say Angular, I mean the modern version, which is on version 19 now. The old AngularJS is dead. They stopped supporting it in 2021.

The big thing Angular is built for is single-page applications. So instead of your browser reloading the whole page every time you click a link, Angular just swaps out the parts that changed. It makes everything feel fast, more like a native app than a traditional website.""")


# ─── SLIDE 4: Section Header - Use Case ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "Key Compelling Use Case")
add_notes(slide, "OK so now you know what Angular is. But why would you actually pick it over something else?")


# ─── SLIDE 5: Use Case ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Built for Large-Scale Applications",
             font_size=36, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 1.3, 5.5, 3.0, [
    "TypeScript catches bugs early in large codebases",
    "Opinionated structure: every Angular project looks similar",
    "New team members can ramp up faster",
    "Built-in dependency injection for testing and modularity",
    "Angular CLI generates consistent boilerplate",
], font_size=20)
add_text_box(slide, 0.8, 4.5, 5.5, 0.5, "Who uses it?",
             font_size=24, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 5.1, 11, 2.0, [
    "Google (Cloud Console, Firebase Console)  |  Microsoft (Office Online, Xbox)",
    "Deutsche Bank  |  Forbes  |  Upwork  |  Samsung",
    "Especially popular in enterprise and financial applications",
], font_size=18)
add_notes(slide, """Angular really shines when things get big. Like, if you've got a team of 20 developers all working on the same app, you need some guardrails. And that's exactly what Angular gives you.

Because it's written in TypeScript, you catch a lot of bugs before anyone even runs the code. The compiler just tells you "hey, this is wrong." That's huge when you have a lot of people committing code every day.

And Angular is really opinionated about how you organize your project. Every Angular app looks pretty much the same. So if someone new joins the team, they don't have to spend two weeks figuring out how the project is set up. They already know where things are.

There's also this thing called dependency injection baked in, which I'll show you later. It basically makes testing way easier.

And these aren't just theoretical benefits. Google uses Angular for their Cloud Console, Firebase Console, a bunch of internal tools. Microsoft uses it for Office Online and parts of Xbox. Banks like Deutsche Bank use it. Forbes, Upwork, Samsung. It's really popular in enterprise, especially finance, where reliability and maintainability matter more than being trendy.""")


# ─── SLIDE 6: Section Header - Demo ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "Demo App Walkthrough")
add_notes(slide, "Alright, enough talking about it. Let me actually show you what it looks like to build something. We're going to walk through a simple to-do list app.")


# ─── SLIDE 7: Project Setup ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Setting Up a Project",
             font_size=36, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 1.3, 5.0, 2.0, [
    "Angular CLI creates the full project scaffold",
    "TypeScript config, testing, dev server: all set up automatically",
    "One command to start developing",
], font_size=20)
add_code_box(slide, 6.5, 1.3, 6.0, 2.5, """# Create a new Angular project
ng new todo-app

# Navigate into the project
cd todo-app

# Start the dev server
ng serve

# App is live at http://localhost:4200""", title="TERMINAL")

add_text_box(slide, 0.8, 3.8, 11, 0.5, "What the CLI generates:",
             font_size=20, color=WHITE, bold=True)
add_code_box(slide, 0.8, 4.4, 11.7, 2.8, """todo-app/
  src/
    app/
      app.component.ts      # Root component (TypeScript class)
      app.component.html    # Root template (HTML)
      app.component.css     # Root styles
      app.component.spec.ts # Root test file
      app.routes.ts         # Routing configuration
    index.html              # Entry point
    main.ts                 # Bootstrap file
  angular.json              # Project configuration
  package.json              # Dependencies
  tsconfig.json             # TypeScript configuration""", title="PROJECT STRUCTURE")
add_notes(slide, """So the first thing you do is run "ng new todo-app" in your terminal. That's it. One command and Angular's CLI sets up your entire project. TypeScript config, testing framework, dev server, all the dependencies. You don't have to configure anything.

Then you just type "ng serve" and your app is running at localhost:4200.

And look at what it generates for you down here. You get this clean folder structure with your root component already broken into separate files: one for the logic, one for the HTML, one for the styles, one for tests. Plus all your config files. This is something I keep coming back to with Angular. It just makes these decisions for you. You're not spending an hour googling "what's the best folder structure for a web app." You just start building.""")


# ─── SLIDE 8: Creating a Component ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Creating a Component",
             font_size=36, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 1.3, 5.5, 1.5, [
    "Components are the building blocks of an Angular app",
    "Each component owns a piece of the UI",
    "CLI generates all the files you need with one command",
], font_size=20)

add_code_box(slide, 0.8, 3.0, 5.5, 0.7, """ng generate component todo-list""", title="TERMINAL")

add_code_box(slide, 0.8, 4.0, 5.5, 3.2, """// todo-list.component.ts
import { Component } from '@angular/core';

@Component({
  selector: 'app-todo-list',
  templateUrl: './todo-list.component.html',
  styleUrls: ['./todo-list.component.css']
})
export class TodoListComponent {
  todos: string[] = [];
  newTodo: string = '';
}""", title="COMPONENT CLASS")

add_code_box(slide, 6.8, 3.0, 5.7, 4.2, """<!-- todo-list.component.html -->
<div class="todo-container">
  <h2>My To-Do List</h2>

  <div class="input-row">
    <input [(ngModel)]="newTodo"
           placeholder="Add a task..." />
    <button (click)="addTodo()">Add</button>
  </div>

  <ul>
    <li *ngFor="let todo of todos">
      {{ todo }}
    </li>
  </ul>
</div>""", title="TEMPLATE")
add_notes(slide, """OK so now we want to actually add our to-do list feature. And in Angular, everything is a component. A component is basically a self-contained piece of the UI.

You run "ng generate component todo-list" and it creates four files for you automatically.

So on the left here, this is the TypeScript class. See that @Component decorator at the top? That's how Angular knows this is a component and where to find its template and styles. The class itself is pretty simple. We've got an array for our todos and a string to track what the user's typing.

And on the right is the HTML template. Now this looks like regular HTML but with some Angular-specific stuff mixed in. That [(ngModel)] on the input? That's called two-way data binding, and I'll explain why it's so cool on the next slide. The *ngFor is basically a for loop in HTML. It goes through the todos array and makes a list item for each one. And the double curly braces just print out the value.""")


# ─── SLIDE 9: Data Binding ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Data Binding in Action",
             font_size=36, color=WHITE, bold=True)

add_code_box(slide, 0.8, 1.3, 5.5, 2.5, """<!-- One-way: display data -->
<h1>{{ title }}</h1>
<img [src]="imageUrl" />

<!-- Event binding: handle clicks -->
<button (click)="addTodo()">Add</button>

<!-- Two-way: input stays in sync -->
<input [(ngModel)]="newTodo" />""", title="BINDING SYNTAX")

add_code_box(slide, 6.8, 1.3, 5.7, 2.5, """// Component class
export class TodoListComponent {
  title = 'My Todos';
  newTodo = '';
  todos: string[] = [];

  addTodo() {
    if (this.newTodo.trim()) {
      this.todos.push(this.newTodo);
      this.newTodo = '';  // clears the input
    }
  }
}""", title="THE CODE BEHIND IT")

add_text_box(slide, 0.8, 4.2, 11, 0.5, "Why this matters:",
             font_size=24, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 4.8, 11, 2.5, [
    "No manual DOM manipulation (no document.getElementById, no innerHTML)",
    "The template and the data stay in sync automatically",
    "Two-way binding with [(ngModel)] means typing in the input updates the variable AND vice versa",
    "This is a big productivity win over vanilla JavaScript",
], font_size=18)
add_notes(slide, """OK so data binding is kind of the magic of Angular. Let me walk through the different types.

The simplest one is those double curly braces. You just write {{ title }} and whatever value title has in your code shows up on screen. If it changes, the screen updates. You don't have to do anything.

Square brackets are for binding HTML attributes. So if you write [src]="imageUrl", Angular keeps the image source in sync with your variable.

Parentheses are for events. When someone clicks that button, Angular calls the addTodo function. Pretty straightforward.

But the really cool one is the two-way binding with [(ngModel)]. People call this the "banana in a box" syntax because of the bracket-parenthesis combo. What it does is keep the input field and your variable perfectly in sync. You type something, the variable updates. You change the variable in code, the input updates. It just works both ways.

And the reason this matters is you never have to write document.getElementById or mess with innerHTML. If you've ever done that in vanilla JavaScript, you know how tedious it gets. Angular just handles all of that for you.""")


# ─── SLIDE 10: Event Handling ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Handling User Interaction",
             font_size=36, color=WHITE, bold=True)

add_code_box(slide, 0.8, 1.3, 5.5, 5.5, """// Full component with interactions
export class TodoListComponent {
  todos: {text: string, done: boolean}[] = [];
  newTodo = '';

  addTodo() {
    if (this.newTodo.trim()) {
      this.todos.push({
        text: this.newTodo,
        done: false
      });
      this.newTodo = '';
    }
  }

  toggleTodo(todo) {
    todo.done = !todo.done;
  }

  deleteTodo(index: number) {
    this.todos.splice(index, 1);
  }
}""", title="COMPONENT CLASS")

add_code_box(slide, 6.8, 1.3, 5.7, 5.5, """<!-- Template with interactions -->
<div class="todo-container">
  <h2>My To-Do List</h2>

  <div class="input-row">
    <input [(ngModel)]="newTodo"
           (keyup.enter)="addTodo()"
           placeholder="Add a task..." />
    <button (click)="addTodo()">Add</button>
  </div>

  <ul>
    <li *ngFor="let todo of todos;
                let i = index"
        [class.completed]="todo.done">

      <input type="checkbox"
             [checked]="todo.done"
             (change)="toggleTodo(todo)" />
      {{ todo.text }}
      <button (click)="deleteTodo(i)">X</button>
    </li>
  </ul>
</div>""", title="TEMPLATE")
add_notes(slide, """Now here's the more complete version with all the interactions wired up.

On the left, I've upgraded the todos from plain strings to objects. Each one has a text and a done flag. And we've got three methods now. addTodo pushes a new item onto the list. toggleTodo flips the checkbox. deleteTodo removes one by its position in the array.

On the right, the template has gotten more interesting. See that (keyup.enter)? That means pressing Enter in the input also adds a todo, not just clicking the button. Little quality-of-life thing.

Each list item has a checkbox now, and look at how clean this is. [checked] binds the checkbox to the done flag, and (change) calls toggleTodo when you click it. That [class.completed] part adds a CSS class when the item is done, so you can strike through the text with just CSS. And the X button calls deleteTodo with the index.

What I really like about this is how readable it is. You can basically read the template like English. For each todo, show a checkbox, show the text, show a delete button. Everything about the interaction is right there in the HTML.""")


# ─── SLIDE 11: Section Header - Architecture ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "Design & Architecture Overview")
add_notes(slide, "So we've seen what it looks like to build something. Now let's zoom out and look at the bigger picture of how Angular apps are put together.")


# ─── SLIDE 12: Architecture Overview ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "How Angular Apps Are Structured",
             font_size=36, color=WHITE, bold=True)

# Key abstractions
add_text_box(slide, 0.8, 1.3, 5.5, 0.5, "Key Abstractions",
             font_size=24, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 1.9, 5.5, 4.5, [
    "Components: UI building blocks (a class + template + styles)",
    "Templates: HTML with Angular syntax for binding and logic",
    "Directives: Extend HTML (*ngIf, *ngFor, custom behaviors)",
    "Pipes: Transform displayed data (date, currency, custom)",
    "Services: Shared logic (API calls, state, business rules)",
    "Modules: Group related features together",
    "Router: Maps URLs to components, handles navigation",
], font_size=18)

add_code_box(slide, 6.8, 1.3, 5.7, 5.5, """How the pieces connect:

  Browser URL
      |
  Router  -->  selects a Component
      |
  Component
      |--- Template (HTML view)
      |--- Styles (CSS)
      |--- Class (TypeScript logic)
      |       |
      |       |--- injects Services
      |               |
      |               |--- HTTP calls
      |               |--- shared state
      |               |--- business logic
      |
      |--- Child Components
              |--- (same structure)""", title="ARCHITECTURE DIAGRAM")
add_notes(slide, """So Angular has about seven main building blocks, and they all fit together in a pretty logical way.

Components are the big one. We've already seen these. Every piece of your UI is a component with its own class, template, and styles. Templates are the HTML side with Angular's binding syntax mixed in. Directives are how you extend HTML to do more, like *ngIf to conditionally show something or *ngFor for loops. Pipes are for formatting. You've got a raw date and you want it to look nice? Pipe it through Angular's date pipe and you're done.

Services are where you put logic that doesn't belong in a component. Like API calls, or shared state between different parts of your app. Modules group related features together. And the Router handles navigation, mapping URLs to the right component.

On the right I've drawn out how these pieces connect. A URL comes in, the router picks the right component, that component has its template and styles, the logic injects whatever services it needs, and components can nest inside each other. It's a tree structure all the way down.

And this separation is really what makes Angular apps maintainable. Everything has a place.""")


# ─── SLIDE 13: Services & DI ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Services & Dependency Injection",
             font_size=36, color=WHITE, bold=True)

add_code_box(slide, 0.8, 1.3, 5.5, 3.5, """// todo.service.ts
import { Injectable } from '@angular/core';
import { HttpClient } from '@angular/common/http';

@Injectable({
  providedIn: 'root'  // available app-wide
})
export class TodoService {
  private apiUrl = '/api/todos';

  constructor(private http: HttpClient) {}

  getTodos() {
    return this.http.get(this.apiUrl);
  }

  addTodo(text: string) {
    return this.http.post(this.apiUrl, { text });
  }
}""", title="SERVICE")

add_code_box(slide, 6.8, 1.3, 5.7, 3.5, """// todo-list.component.ts
import { TodoService } from './todo.service';

export class TodoListComponent {
  todos = [];

  // Angular automatically provides
  // the TodoService instance
  constructor(private todoService: TodoService) {}

  ngOnInit() {
    this.todoService.getTodos()
      .subscribe(data => this.todos = data);
  }

  addTodo(text: string) {
    this.todoService.addTodo(text)
      .subscribe(todo => this.todos.push(todo));
  }
}""", title="COMPONENT USING THE SERVICE")

add_text_box(slide, 0.8, 5.2, 11, 0.5, "Why dependency injection matters:",
             font_size=24, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 5.8, 11, 1.5, [
    "Components don't create their own dependencies, Angular provides them",
    "Easy to swap implementations (real API vs mock for testing)",
    "Services are singletons by default, so shared state just works",
], font_size=18)
add_notes(slide, """This is the part that really sets Angular apart from most other front-end tools. Dependency injection. If you've used Java or C#, you've probably seen this pattern before, but it's pretty rare in front-end frameworks.

So on the left, we've got a service. Think of it as a helper class that handles stuff your component shouldn't worry about, like talking to an API. That @Injectable decorator at the top is how Angular knows it can inject this into other things. And providedIn root means it's available everywhere in the app.

Now look at the component on the right. See the constructor? It just says "hey, I need a TodoService." That's it. The component never creates the service itself. Angular sees that constructor parameter and automatically provides an instance. The component just says what it needs and Angular handles the wiring.

Why does this matter? Three reasons. First, your components stay clean. They just deal with the UI, not how data gets fetched. Second, testing gets way easier. You can swap in a fake service without changing any component code. And third, services are shared by default. If two different components inject the same service, they get the same instance. So sharing data between components is basically free.""")


# ─── SLIDE 14: Section Header - Status ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "Status & Adoption")
add_notes(slide, "Alright, let's talk about where Angular stands right now.")


# ─── SLIDE 15: Status ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Status & Adoption",
             font_size=36, color=WHITE, bold=True)

# Stats in a grid-like layout
stats = [
    ("Open Source", "MIT License\nMaintained by Google"),
    ("Current Version", "Angular 19\nReleased Nov 2024"),
    ("Release Cycle", "Major version every 6 months\n18 months support per version"),
    ("npm Downloads", "~3-4 million\nweekly downloads"),
    ("GitHub Stars", "~96,000 stars\nActive contributor community"),
    ("Ecosystem", "Angular Material, NgRx,\nPrimeNG, and more"),
]
for i, (label, value) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.5 + row * 2.5

    add_text_box(slide, x, y, 3.5, 0.4, label,
                 font_size=16, color=ANGULAR_RED, bold=True)
    add_text_box(slide, x, y + 0.45, 3.5, 1.5, value,
                 font_size=20, color=LIGHT_GRAY)
add_notes(slide, """So Angular is fully open source, MIT license, and it's maintained by a dedicated team at Google. This isn't some side project. Google has engineers working on it full time.

They're on version 19 right now, which came out in November 2024. And they're pretty consistent with releases. New major version every six months, usually May and November. Each version gets a year and a half of support total. Six months of active updates, then twelve months where they only do security patches.

Community-wise, it's big. Around 3 to 4 million downloads per week on npm. About 96,000 stars on GitHub. And there's a whole ecosystem around it. Angular Material gives you pre-built UI components that follow Google's Material Design. NgRx handles complex state management. PrimeNG gives you a huge library of widgets. So you're not on your own.""")


# ─── SLIDE 16: Section Header - Competitors ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "Nearest Competitors")
add_notes(slide, "OK so Angular doesn't exist in a vacuum. Let's see how it compares to the other big players.")


# ─── SLIDE 17: Competitors ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Angular vs the Competition",
             font_size=36, color=WHITE, bold=True)

# Comparison table as text
headers = "                Angular          React            Vue              Svelte"
rows = [
    "Type            Full framework   UI library       Progressive fw   Compiler-based",
    "Language        TypeScript       JS or TS         JS or TS         JS or TS",
    "                (required)       (optional)       (optional)       (optional)",
    "",
    "Tooling         Everything       Bring your own   Middle ground    SvelteKit",
    "                included",
    "",
    "Learning        Steepest         Moderate         Gentlest         Low",
    "Curve",
    "",
    "Bundle Size     Largest          Medium           Smaller          Smallest",
    "",
    "Best For        Large enterprise Flexible,        Quick ramp-up    Performance-",
    "                apps             any scale                         critical apps",
]
add_code_box(slide, 0.8, 1.3, 11.7, 5.8,
             headers + "\n" + "-" * 80 + "\n" + "\n".join(rows),
             title="COMPARISON")
add_notes(slide, """So here's the comparison table. Four main options in the front-end world right now.

React is the most popular by a lot, but here's the thing, it's technically a library, not a framework. It handles the view layer and that's it. You need a router? Go find one. State management? Pick from five options. Forms? Figure it out. That's great if you want flexibility, but it also means a lot of decisions and a lot of glue code.

Vue is what I'd call the friendliest option. Easiest learning curve, good documentation, and you can adopt it gradually. Start with a little Vue on one page and scale up. The ecosystem is smaller than React or Angular though.

Svelte is the interesting newcomer. It works completely differently. Instead of running a framework in the browser, it compiles your code into vanilla JavaScript at build time. So there's basically no framework overhead at runtime, which gives you the smallest possible bundle sizes. But the community is still the smallest of the four.

And then Angular is on the far end of the spectrum. It requires TypeScript, comes with everything, and yeah, it has the steepest learning curve. But if you're on a big team building something complex, that structure is exactly what you want.""")


# ─── SLIDE 18: Section Header - When to Use ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 7, "When to Use It / When Not To")
add_notes(slide, "So let's bring it all together. When should you actually reach for Angular?")


# ─── SLIDE 19: When to Use ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SLIDE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "When Should You Use Angular?",
             font_size=36, color=WHITE, bold=True)

add_text_box(slide, 0.8, 1.4, 5.5, 0.5, "Use Angular when:",
             font_size=24, color=GREEN, bold=True)
add_bullet_list(slide, 0.8, 2.0, 5.5, 3.0, [
    "Building a large, long-lived application",
    "Working with a big team that needs consistency",
    "You want everything included out of the box",
    "Your team knows TypeScript or wants enforced type safety",
    "You need enterprise features (DI, strong testing, structure)",
], font_size=18)

add_text_box(slide, 6.8, 1.4, 5.5, 0.5, "Skip Angular when:",
             font_size=24, color=ANGULAR_RED, bold=True)
add_bullet_list(slide, 6.8, 2.0, 5.5, 3.0, [
    "Building a small or simple website (it's overkill)",
    "You need to prototype fast (learning curve + boilerplate)",
    "Your team wants flexibility in choosing libraries",
    "Bundle size is a top priority (Svelte or Vue are lighter)",
], font_size=18)

# Summary box
shape = slide.shapes.add_shape(
    1, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = CODE_BG
shape.line.fill.background()

add_text_box(slide, 1.2, 5.2, 11, 0.5, "One-liner:",
             font_size=20, color=ANGULAR_RED, bold=True)
add_text_box(slide, 1.2, 5.7, 11, 0.8,
             'Angular is the "batteries-included" framework. Pick it when structure and '
             'scalability matter more than speed and simplicity.',
             font_size=22, color=WHITE)
add_notes(slide, """On the left, the green side, these are the situations where Angular is your best bet. You're building something big that's going to be around for years. You've got a large team and you need consistency. You don't want to waste time assembling your own toolkit from scratch. Your team already knows TypeScript or you want that type safety enforced. Or you need those enterprise features like dependency injection and solid testing infrastructure.

On the right, the red side, this is where you should probably look elsewhere. If you're just making a landing page or a simple website, Angular is like bringing a semi truck to pick up groceries. If you need to prototype something fast, the learning curve and all that boilerplate are going to slow you down. If your team likes picking their own tools, Angular's going to feel like a straitjacket. And if you really care about bundle size, Svelte and Vue will be lighter.

And the one-liner I want to leave you with is at the bottom. Angular is the batteries-included framework. Pick it when structure and scalability matter more than speed and simplicity.""")


# ─── SLIDE 20: Thank You ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 2.5, 11, 1.0, "Thank You",
             font_size=54, color=WHITE, bold=True)
shape = slide.shapes.add_shape(1, Inches(0.8), Inches(3.6), Inches(3), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ANGULAR_RED
shape.line.fill.background()
add_text_box(slide, 0.8, 4.0, 11, 0.5, "Questions?",
             font_size=28, color=LIGHT_GRAY)
add_text_box(slide, 0.8, 6.5, 11, 0.4,
             "AI Disclosure: Outline and research notes prepared with assistance from Claude (Anthropic).",
             font_size=12, color=MID_GRAY)
add_notes(slide, "And that's Angular. Thanks for watching, and I'm happy to take any questions.")


# Save
output_path = r"C:\Users\lawss\Documents\Repos\User Interface Design\modules\module-11\Angular-Tech-Talk.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
